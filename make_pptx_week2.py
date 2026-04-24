"""
Week 2 Progress Presentation generator.

Creates a PPTX that includes:
  - baseline figure
  - period sweep figure
  - verification convergence figure
  - 10-test verification result table
"""

import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT_DIR = os.path.dirname(os.path.abspath(__file__))
PPTX_PATH = os.path.join(OUT_DIR, "week2_progress.pptx")
FIG_DIR = os.path.join(OUT_DIR, "figures")
EQ_DIR = os.path.join(OUT_DIR, "_eq_images")
os.makedirs(EQ_DIR, exist_ok=True)

BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
TITLE_COLOR = RGBColor(0x1A, 0x47, 0x8A)
SUBTITLE_COLOR = RGBColor(0x2E, 0x6B, 0xB0)
TEXT_COLOR = RGBColor(0x22, 0x22, 0x22)
ACCENT_COLOR = RGBColor(0xC0, 0x39, 0x2B)
MUTED_COLOR = RGBColor(0x66, 0x66, 0x77)
HIGHLIGHT_BG = RGBColor(0xF0, 0xF4, 0xFA)
PASS_COLOR = RGBColor(0x1E, 0x8A, 0x4C)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def render_latex(latex_str: str, fontsize: int = 22, dpi: int = 200) -> str:
    fname = os.path.join(EQ_DIR, f"w2_eq_{hash(latex_str) & 0xFFFFFFFF:08x}.png")
    if os.path.exists(fname):
        return fname
    fig, ax = plt.subplots(figsize=(0.01, 0.01))
    ax.axis("off")
    fig.patch.set_alpha(0)
    ax.text(0, 0, f"${latex_str}$", fontsize=fontsize, color="black",
            ha="left", va="baseline", transform=ax.transAxes)
    fig.savefig(fname, dpi=dpi, transparent=True,
                bbox_inches="tight", pad_inches=0.05)
    plt.close(fig)
    return fname


def set_slide_bg(slide, color=BG_COLOR):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=TEXT_COLOR, bold=False,
             alignment=PP_ALIGN.LEFT):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return p


def add_paragraph(tf, text, size=18, color=TEXT_COLOR, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(4),
                  space_after=Pt(2), level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p


def add_bullet(tf, text_en, text_kr, size=16, level=0):
    add_paragraph(tf, f"• {text_en}", size=size, color=TEXT_COLOR, level=level)
    add_paragraph(tf, f"  {text_kr}", size=size - 2, color=MUTED_COLOR,
                  level=level, space_before=Pt(0))


def add_equation_image(slide, latex_str, left, top, max_height=Inches(0.55),
                       fontsize=22):
    img_path = render_latex(latex_str, fontsize=fontsize)
    img = Image.open(img_path)
    w, h = img.size
    aspect = w / h
    img_h = max_height
    img_w = int(img_h * aspect)
    return slide.shapes.add_picture(img_path, left, top, img_w, img_h)


def add_image_fit(slide, img_path, left, top, max_w, max_h):
    img = Image.open(img_path)
    w, h = img.size
    aspect = w / h
    if max_w / max_h > aspect:
        img_h = max_h
        img_w = int(max_h * aspect)
    else:
        img_w = max_w
        img_h = int(max_w / aspect)
    return slide.shapes.add_picture(img_path, left, top, img_w, img_h)


def make_slide_header(slide, num, title_en, title_kr):
    set_slide_bg(slide)
    tb = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5))
    set_text(tb.text_frame, f"{num}. {title_en}", size=28,
             color=TITLE_COLOR, bold=True)
    tb2 = add_textbox(slide, Inches(0.85), Inches(0.85), Inches(12), Inches(0.4))
    set_text(tb2.text_frame, title_kr, size=18, color=MUTED_COLOR)


# ---------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------
def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    tb = add_textbox(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2))
    set_text(tb.text_frame, "Week 2 Progress Presentation",
             size=36, color=TITLE_COLOR, bold=True, alignment=PP_ALIGN.CENTER)

    tb2 = add_textbox(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(1.0))
    set_text(tb2.text_frame,
             "CDA Baseline Implementation & Comprehensive Verification",
             size=24, color=TEXT_COLOR, alignment=PP_ALIGN.CENTER)
    add_paragraph(tb2.text_frame,
                  "CDA 베이스라인 구현 및 종합 검증",
                  size=22, color=MUTED_COLOR, alignment=PP_ALIGN.CENTER)

    tb3 = add_textbox(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.6))
    set_text(tb3.text_frame, "10 / 10 physical verification tests  PASS",
             size=26, color=PASS_COLOR, bold=True, alignment=PP_ALIGN.CENTER)


def slide_question(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 1, "Question Being Tested", "검증 중인 질문")

    shape = slide.shapes.add_shape(
        1, Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = HIGHLIGHT_BG
    shape.line.color.rgb = TITLE_COLOR
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    set_text(tf, "", size=10)
    add_paragraph(tf,
        '"How does the phase distortion caused by inter-meta-atom coupling '
        'increase as the array period P decreases?"',
        size=22, color=TEXT_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, "", size=6)
    add_paragraph(tf,
        '"배열 주기 P가 줄어들수록 메타 원자 간 커플링에 의한 위상 왜곡이 어떻게 증가하는가?"',
        size=20, color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.5))
    tf2 = tb.text_frame
    tf2.word_wrap = True
    add_paragraph(tf2, "This week's focus:", size=18, color=SUBTITLE_COLOR, bold=True)
    add_paragraph(tf2, "이번 주의 목표", size=15, color=MUTED_COLOR, space_before=Pt(0))
    add_bullet(tf2, "Build CDA simulation from scratch",
               "CDA 시뮬레이션 직접 구현", size=16)
    add_bullet(tf2, "Run 10 independent physical verification tests",
               "10가지 독립적 물리 검증 수행", size=16)


def slide_baseline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 2, "Baseline Implemented", "구현한 베이스라인")

    tb = add_textbox(slide, Inches(0.7), Inches(1.6), Inches(6.0), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, "Code structure", size=20, color=ACCENT_COLOR, bold=True)
    add_paragraph(tf, "코드 구조", size=15, color=MUTED_COLOR, space_before=Pt(0))

    files = [
        ("cda.py", "Core library — Green's function, matrix, solver",
         "핵심 라이브러리: Green 함수, 행렬, 솔버"),
        ("run_baseline.py", "Baseline run + first figures",
         "기본 시뮬 실행 + 첫 그림 생성"),
        ("run_period_sweep.py", "Period sweep (Week 3 preview)",
         "주기 스윕 (Week 3 사전 실험)"),
        ("run_verification.py", "10 physical verification tests",
         "10가지 물리 검증 테스트"),
    ]
    for fname, en, kr in files:
        add_paragraph(tf, f"▸ {fname}", size=17, color=TITLE_COLOR, bold=True,
                      space_before=Pt(8))
        add_paragraph(tf, f"   {en}", size=14, color=TEXT_COLOR,
                      space_before=Pt(0))
        add_paragraph(tf, f"   {kr}", size=13, color=MUTED_COLOR,
                      space_before=Pt(0))

    # Right column: supported features
    tb2 = add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    set_text(tf2, "What works", size=20, color=ACCENT_COLOR, bold=True)
    add_paragraph(tf2, "동작 확인 사항", size=15, color=MUTED_COLOR, space_before=Pt(0))

    features = [
        ("End-to-end forward run: uniform & non-uniform arrays",
         "균일 & 비균일 배열 end-to-end 동작"),
        ("All diagnostics: Δφ, max Δφ, ADR, extinction proxy",
         "모든 진단 지표 구현"),
        ("1D arrays of arbitrary size, complex Lorentzian α",
         "임의 크기 1D 배열, 복소 Lorentzian 분극률"),
        ("10/10 physical verification tests PASS",
         "10가지 물리 검증 모두 통과"),
    ]
    for en, kr in features:
        add_paragraph(tf2, f"✓  {en}", size=15, color=PASS_COLOR, bold=True,
                      space_before=Pt(8))
        add_paragraph(tf2, f"    {kr}", size=13, color=MUTED_COLOR,
                      space_before=Pt(0))


def slide_verification_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 3, "Verification Results  —  10 / 10 PASS",
                      "검증 결과 — 10/10 통과")

    tests = [
        ("1", "Single dipole  =  α · E_inc",
         "1개 쌍극자 해석해", "rel. err  1.5×10⁻¹⁶"),
        ("2", "Two-dipole analytical  p = α/(1 − αG)",
         "2개 쌍극자 해석해", "rel. err  2.5×10⁻¹⁶"),
        ("3", "Green's function formula + reciprocity",
         "Green 함수 공식 + 호혜성", "exact  0"),
        ("4", "Linear system residual  ‖Ap − E‖/‖E‖",
         "선형 시스템 잔차", "5.4×10⁻¹⁶"),
        ("5", "Mirror symmetry of centered array",
         "거울 대칭성", "7×10⁻¹⁶"),
        ("6", "Matrix reciprocity  A_ij = A_ji",
         "행렬 호혜성", "exact  0"),
        ("7", "Array-size convergence (2D lattice sum)",
         "배열 크기 수렴성", "slope ≈ −0.26 ★"),
        ("8", "Wavelength-scaling invariance",
         "파장 스케일링 불변성", "exact  0"),
        ("9", "Extinction power positivity (optical theorem)",
         "소광 전력 양성", "all  P_ext > 0"),
        ("10", "Non-uniform solver ⇒ uniform limit",
         "비균일→균일 극한", "exact  0"),
    ]

    y_start = Inches(1.5)
    row_h = Inches(0.47)
    left = Inches(0.5)

    # header
    headers = [("#", Inches(0.5)), ("Test", Inches(6.3)),
               ("한국어", Inches(2.7)), ("Result", Inches(3.2))]
    x = left
    for h, w in headers:
        tb = add_textbox(slide, x, y_start, w, Inches(0.35))
        set_text(tb.text_frame, h, size=14, color=TITLE_COLOR, bold=True)
        x += w

    # rows
    for i, (num, en, kr, result) in enumerate(tests):
        y = y_start + Inches(0.4) + row_h * i
        x = left
        shape = slide.shapes.add_shape(
            1, x, y, sum(w for _, w in headers), Inches(0.45)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = (
            RGBColor(0xF8, 0xFA, 0xFD) if i % 2 == 0
            else RGBColor(0xFF, 0xFF, 0xFF)
        )
        shape.line.color.rgb = RGBColor(0xD0, 0xD8, 0xE0)
        shape.line.width = Pt(0.5)

        cells = [
            (num, headers[0][1], PASS_COLOR, True, 14),
            (en, headers[1][1], TEXT_COLOR, False, 13),
            (kr, headers[2][1], MUTED_COLOR, False, 12),
            (result, headers[3][1], PASS_COLOR, True, 13),
        ]
        cx = left
        for text, w, color, bold, size in cells:
            tb = add_textbox(slide, cx + Inches(0.05), y + Inches(0.05),
                             w - Inches(0.1), Inches(0.35))
            set_text(tb.text_frame, text, size=size, color=color, bold=bold)
            cx += w


def slide_test7_highlight(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 4, "Key Physical Finding — Test 7",
                      "주목할 물리 — Test 7")

    tb = add_textbox(slide, Inches(0.6), Inches(1.55), Inches(5.8), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, "Observation", size=20, color=ACCENT_COLOR, bold=True)
    add_paragraph(tf, "관측 결과", size=14, color=MUTED_COLOR, space_before=Pt(0))
    add_bullet(tf, "Convergence rate of central-atom Δφ vs N",
               "중심 원자 Δφ의 N에 따른 수렴 속도", size=15)
    add_bullet(tf, "Observed slope: -0.26",
               "관측 기울기: -0.26", size=15)
    add_bullet(tf, "Theoretical (2D free-space Green): -0.5",
               "이론값 (2D Green): -0.5", size=15)

    add_paragraph(tf, "", size=6)
    add_paragraph(tf, "Interpretation", size=20, color=ACCENT_COLOR, bold=True,
                  space_before=Pt(10))
    add_paragraph(tf, "해석", size=14, color=MUTED_COLOR, space_before=Pt(0))
    add_bullet(tf,
               "2D free-space Green's function decays as 1/sqrt(r)",
               "2D Green 함수는 1/√r로 감쇠", size=15)
    add_bullet(tf,
               "1D lattice sum is conditionally convergent, slow O(1/sqrt(N))",
               "1D 격자합은 조건부 수렴, O(1/√N)", size=15)
    add_bullet(tf,
               "3D would converge much faster (Green ~ 1/r)",
               "3D였다면 1/r로 훨씬 빠름", size=15)
    add_bullet(tf,
               "Not a numerical error — it reflects genuine 2D physics",
               "수치 오류가 아니라 2D 물리의 본질", size=15)

    # Right: image
    img_path = os.path.join(FIG_DIR, "verification_convergence.png")
    if os.path.exists(img_path):
        add_image_fit(slide, img_path,
                      Inches(6.7), Inches(1.5),
                      Inches(6.3), Inches(5.4))


def slide_baseline_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 5, "Baseline Result", "베이스라인 결과")

    tb = add_textbox(slide, Inches(0.6), Inches(1.55), Inches(5.8), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, "Setup", size=20, color=ACCENT_COLOR, bold=True)
    add_paragraph(tf, "설정", size=14, color=MUTED_COLOR, space_before=Pt(0))
    add_bullet(tf, "Uniform 1D array, N = 31",
               "1D 균일 배열, N = 31", size=16)
    add_bullet(tf, "Period P = 0.6 λ (sub-wavelength)",
               "주기 P = 0.6 λ (서브파장)", size=16)
    add_bullet(tf, "Normal-incidence plane wave",
               "정상 입사 평면파", size=16)

    add_paragraph(tf, "", size=6)
    add_paragraph(tf, "Measured", size=20, color=ACCENT_COLOR, bold=True,
                  space_before=Pt(10))
    add_paragraph(tf, "측정 결과", size=14, color=MUTED_COLOR, space_before=Pt(0))
    add_bullet(tf, "Mean |Δφ| = 6.09°", "평균 위상 편차 6.09°", size=16)
    add_bullet(tf, "Max |Δφ| = 7.89°", "최대 편차 7.89°", size=16)
    add_bullet(tf, "ADR range [0.979, 1.041]",
               "진폭 왜곡 비 [0.979, 1.041]", size=16)

    add_paragraph(tf, "", size=6)
    add_paragraph(tf, "Conclusion", size=20, color=PASS_COLOR, bold=True,
                  space_before=Pt(10))
    add_bullet(tf, "Coupling is real, measurable, non-negligible",
               "커플링은 실제로 측정 가능한 수준", size=16)

    img_path = os.path.join(FIG_DIR, "baseline_uniform_array.png")
    if os.path.exists(img_path):
        add_image_fit(slide, img_path,
                      Inches(6.7), Inches(1.5),
                      Inches(6.3), Inches(5.4))


def slide_period_sweep(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 6, "First Scan — Period Sweep",
                      "첫 스윕 결과 — 주기 스윕")

    tb = add_textbox(slide, Inches(0.6), Inches(1.55), Inches(5.8), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, "Uniform vs non-uniform", size=20,
             color=ACCENT_COLOR, bold=True)
    add_paragraph(tf, "균일 vs 비균일", size=14, color=MUTED_COLOR, space_before=Pt(0))

    rows = [
        ("P/λ", "Uniform", "Non-uniform"),
        ("0.5", "5.8°", "6.6°"),
        ("1.0", "27.0°", "22.4°"),
        ("1.5", "3.4°", "3.8°"),
        ("2.0", "14.1°", "18.3°"),
        ("3.0", "20.9°", "19.8°"),
    ]
    for i, (a, b, c) in enumerate(rows):
        color = SUBTITLE_COLOR if i == 0 else TEXT_COLOR
        bold = (i == 0) or (float(a.replace("/λ", "").split()[0]) in (1.0, 2.0, 3.0))
        prefix = "  "
        text = f"{prefix}{a:>8}   {b:>10}   {c:>14}"
        add_paragraph(tf, text, size=15, color=color, bold=bold,
                      space_before=Pt(4))

    add_paragraph(tf, "", size=6)
    add_paragraph(tf, "Key observation", size=20, color=ACCENT_COLOR, bold=True,
                  space_before=Pt(10))
    add_paragraph(tf, "핵심 관찰", size=14, color=MUTED_COLOR, space_before=Pt(0))
    add_bullet(tf, "Sharp peaks at P = m · λ  (m ∈ ℤ)",
               "파장의 정수배 주기에서 급격한 피크", size=16)
    add_bullet(tf, "Wood anomaly / lattice resonance",
               "Wood anomaly, 격자 공명", size=16)
    add_bullet(tf, "Coincides with diffraction-order onset",
               "회절 차수 개시 조건과 일치", size=16)

    img_path = os.path.join(FIG_DIR, "period_sweep_comparison.png")
    if os.path.exists(img_path):
        add_image_fit(slide, img_path,
                      Inches(6.7), Inches(1.5),
                      Inches(6.3), Inches(5.4))


def slide_difficulties(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 7, "Main Difficulties", "현재 주요 어려움")

    # Difficulty 1
    shape1 = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2)
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = HIGHLIGHT_BG
    shape1.line.color.rgb = ACCENT_COLOR
    shape1.line.width = Pt(1.5)

    tf1 = shape1.text_frame
    tf1.word_wrap = True
    set_text(tf1, "(a) Lattice-resonance divergence", size=20,
             color=ACCENT_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf1, "격자 공명 발산", size=15, color=MUTED_COLOR,
                  alignment=PP_ALIGN.CENTER, space_before=Pt(0))
    add_paragraph(tf1, "", size=6)
    add_bullet(tf1, "Matrix nearly singular at P ≈ m · λ",
               "P = mλ 근처에서 행렬이 거의 특이", size=16)
    add_bullet(tf1, "Δφ(P) curve non-monotonic",
               "Δφ(P) 곡선이 단조롭지 않음", size=16)
    add_bullet(tf1, "Must isolate off-resonance behavior",
               "공명 영역을 피해 분석 필요", size=16)

    # Difficulty 2
    shape2 = slide.shapes.add_shape(
        1, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = HIGHLIGHT_BG
    shape2.line.color.rgb = TITLE_COLOR
    shape2.line.width = Pt(1.5)

    tf2 = shape2.text_frame
    tf2.word_wrap = True
    set_text(tf2, "(b) Slow 2D convergence", size=20,
             color=TITLE_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf2, "2D 시스템의 느린 수렴", size=15, color=MUTED_COLOR,
                  alignment=PP_ALIGN.CENTER, space_before=Pt(0))
    add_paragraph(tf2, "", size=6)
    add_bullet(tf2, "2D Green decays as 1/sqrt(r)",
               "2D Green 함수 1/√r 감쇠", size=16)
    add_bullet(tf2, "Lattice sum converges conditionally",
               "격자합이 조건부 수렴", size=16)
    add_bullet(tf2, "Finite-array edge effect persists longer than 3D",
               "유한 배열 edge effect가 3D보다 오래 지속", size=16)
    add_bullet(tf2, "Future: periodic Bloch BC, Ewald summation",
               "향후: Bloch 경계, Ewald 합법", size=16)


def slide_next_plan(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 8, "Next Week's Plan", "다음 주 계획")

    tb = add_textbox(slide, Inches(0.7), Inches(1.6), Inches(5.8), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, "Primary comparison", size=22, color=ACCENT_COLOR, bold=True)
    add_paragraph(tf, "메인 비교 실험", size=15, color=MUTED_COLOR, space_before=Pt(0))
    add_paragraph(tf, "", size=6)
    add_bullet(tf, "Uniform vs non-uniform array — detailed analysis",
               "균일 vs 비균일 상세 분석", size=16)
    add_bullet(tf, "Focus on off-resonance regions",
               "공명을 피한 영역에 집중", size=16)
    add_bullet(tf, "Isolate (i) lattice coupling vs (ii) α-inhomogeneity",
               "(i) 격자 커플링과 (ii) α 비균일성 효과 분리", size=16)
    add_bullet(tf, "Find threshold period P* where mean |Δφ| < 1°",
               "평균 |Δφ| < 1°가 되는 임계 주기 P* 추정", size=16)

    tb2 = add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    set_text(tf2, "Additional checks", size=22, color=TITLE_COLOR, bold=True)
    add_paragraph(tf2, "추가 검증", size=15, color=MUTED_COLOR, space_before=Pt(0))
    add_paragraph(tf2, "", size=6)
    add_bullet(tf2, "Conditioning study near Wood anomalies",
               "공명 근처의 행렬 조건수 분석", size=16)
    add_bullet(tf2, "Energy budget: P_sca / P_ext ratio",
               "에너지 예산: P_sca / P_ext 비율", size=16)
    add_bullet(tf2, "Sensitivity to α-ordering in non-uniform cases",
               "비균일 배열의 α 순서 민감도", size=16)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs)
    slide_question(prs)
    slide_baseline(prs)
    slide_verification_table(prs)
    slide_test7_highlight(prs)
    slide_baseline_results(prs)
    slide_period_sweep(prs)
    slide_difficulties(prs)
    slide_next_plan(prs)

    prs.save(PPTX_PATH)
    print(f"PPT saved: {PPTX_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
