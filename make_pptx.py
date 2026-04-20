"""
Week 1 Proposal PPT Generator
메타 원자 간 커플링이 위상 응답에 미치는 영향: Coupled Dipole Approximation
"""

import os
import io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

OUT_DIR = os.path.dirname(os.path.abspath(__file__))
PPTX_PATH = os.path.join(OUT_DIR, "week1_proposal_v2.pptx")
EQ_DIR = os.path.join(OUT_DIR, "_eq_images")
os.makedirs(EQ_DIR, exist_ok=True)

BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
TITLE_COLOR = RGBColor(0x1A, 0x47, 0x8A)
SUBTITLE_COLOR = RGBColor(0x2E, 0x6B, 0xB0)
TEXT_COLOR = RGBColor(0x22, 0x22, 0x22)
ACCENT_COLOR = RGBColor(0xC0, 0x39, 0x2B)
MUTED_COLOR = RGBColor(0x66, 0x66, 0x77)
HIGHLIGHT_BG = RGBColor(0xF0, 0xF4, 0xFA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def render_latex(latex_str: str, fontsize: int = 22, dpi: int = 200) -> str:
    """Render LaTeX string to PNG image, return file path."""
    safe = latex_str.replace("\\", "_bs_").replace("{", "").replace("}", "")
    safe = safe.replace("/", "").replace("$", "").replace(" ", "")[:60]
    fname = os.path.join(EQ_DIR, f"eq_{hash(latex_str) & 0xFFFFFFFF:08x}.png")
    if os.path.exists(fname):
        return fname

    fig, ax = plt.subplots(figsize=(0.01, 0.01))
    ax.axis("off")
    fig.patch.set_alpha(0)
    text = ax.text(
        0, 0, f"${latex_str}$",
        fontsize=fontsize, color="black",
        ha="left", va="baseline",
        transform=ax.transAxes,
    )
    fig.savefig(
        fname, dpi=dpi, transparent=True,
        bbox_inches="tight", pad_inches=0.05,
    )
    plt.close(fig)
    return fname


def set_slide_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=TEXT_COLOR, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return p


def add_paragraph(tf, text, size=18, color=TEXT_COLOR, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(2),
                  level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p


def add_bullet(tf, text_en, text_kr, size=16, level=0):
    add_paragraph(tf, f"• {text_en}", size=size, color=TEXT_COLOR, level=level)
    add_paragraph(tf, f"  {text_kr}", size=size - 2, color=MUTED_COLOR, level=level,
                  space_before=Pt(0))


def add_equation_image(slide, latex_str, left, top, max_height=Inches(0.55),
                       fontsize=22):
    img_path = render_latex(latex_str, fontsize=fontsize)
    img = Image.open(img_path)
    w, h = img.size
    aspect = w / h
    img_h = max_height
    img_w = int(img_h * aspect)
    pic = slide.shapes.add_picture(img_path, left, top, img_w, img_h)
    return pic


def make_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    tb = add_textbox(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1.2))
    set_text(tb.text_frame, "Week 1 Proposal Presentation",
             size=36, color=TITLE_COLOR, bold=True, alignment=PP_ALIGN.CENTER)

    tb2 = add_textbox(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(1.0))
    set_text(tb2.text_frame,
             "Effect of Inter-Meta-Atom Coupling on Phase Response",
             size=24, color=TEXT_COLOR, alignment=PP_ALIGN.CENTER)
    add_paragraph(tb2.text_frame,
                  "메타 원자 간 커플링이 위상 응답에 미치는 영향",
                  size=22, color=MUTED_COLOR, alignment=PP_ALIGN.CENTER)

    tb3 = add_textbox(slide, Inches(1), Inches(4.8), Inches(11.3), Inches(0.8))
    set_text(tb3.text_frame,
             "Coupled Dipole Approximation (CDA)",
             size=28, color=ACCENT_COLOR, bold=True, alignment=PP_ALIGN.CENTER)

    tb4 = add_textbox(slide, Inches(1), Inches(6.0), Inches(11.3), Inches(0.5))
    set_text(tb4.text_frame, "Track A — Build a New Simulation",
             size=18, color=MUTED_COLOR, alignment=PP_ALIGN.CENTER)


def make_slide_header(slide, num, title_en, title_kr):
    set_slide_bg(slide)
    tb = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5))
    set_text(tb.text_frame, f"{num}. {title_en}", size=28, color=TITLE_COLOR, bold=True)

    tb2 = add_textbox(slide, Inches(0.85), Inches(0.85), Inches(12), Inches(0.4))
    set_text(tb2.text_frame, title_kr, size=18, color=MUTED_COLOR)


def slide_1_context(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 1, "Research Context", "연구 배경")

    # Left column - Problem
    tb = add_textbox(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True

    set_text(tf, "Assumption in Metasurface Design", size=20, color=ACCENT_COLOR, bold=True)
    add_paragraph(tf, "메타서페이스 설계의 기본 가정", size=16, color=MUTED_COLOR,
                  space_before=Pt(0))

    add_bullet(tf, "Each meta-atom treated as isolated",
               "각 메타 원자를 독립적(isolated)으로 가정")
    add_bullet(tf, "Individual scattering properties computed first, then applied to array",
               "개별 산란 특성을 먼저 계산 후 배열에 적용")

    add_paragraph(tf, "", size=8)
    add_paragraph(tf, "The Problem: Coupling", size=20, color=ACCENT_COLOR, bold=True)
    add_paragraph(tf, "문제: 상호작용(Coupling) 효과", size=16, color=MUTED_COLOR,
                  space_before=Pt(0))

    add_bullet(tf, "Adjacent meta-atoms are EM coupled in reality",
               "실제로는 인접 원자 간 전자기 커플링 존재")
    add_bullet(tf, "Scattered light from one atom alters neighbors' response",
               "한 원자가 산란한 빛이 인접 원자의 응답을 변화")
    add_bullet(tf, "Coupling strengthens as period decreases",
               "주기가 작을수록 커플링 강화")

    # Right column - Why it matters
    tb2 = add_textbox(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    set_text(tf2, "Why It Matters", size=20, color=ACCENT_COLOR, bold=True)
    add_paragraph(tf2, "왜 중요한가?", size=16, color=MUTED_COLOR, space_before=Pt(0))

    add_bullet(tf2, "Phase precision is critical for metalenses, holograms",
               "고효율 메타렌즈, 홀로그램에서 위상 정밀도가 핵심")
    add_bullet(tf2, "Coupling-induced phase distortion degrades performance",
               "커플링으로 인한 위상 왜곡 → 성능 저하")
    add_bullet(tf2, "Quantitative understanding enables better design",
               "정량적 이해 → 더 정확한 설계 가능")

    add_paragraph(tf2, "", size=8)
    add_paragraph(tf2, "Design vs Reality Mismatch", size=20, color=ACCENT_COLOR, bold=True)
    add_paragraph(tf2, "설계 ↔ 실제 성능 불일치", size=16, color=MUTED_COLOR,
                  space_before=Pt(0))

    add_bullet(tf2, "Isolated assumption → incorrect phase at small periods",
               "고립 가정 → 작은 주기에서 부정확한 위상")
    add_bullet(tf2, "Need to quantify: how much distortion, at what period?",
               "얼마나 왜곡되는지, 어떤 주기에서? → 정량화 필요")


def slide_2_question(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 2, "Research Question", "연구 질문")

    # Main question box
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1, Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.0)  # 1 = rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = HIGHLIGHT_BG
    shape.line.color.rgb = TITLE_COLOR
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    set_text(tf, "", size=10)
    add_paragraph(tf,
        '"How does the phase distortion caused by inter-meta-atom coupling '
        'increase as the array period decreases?"',
        size=22, color=TEXT_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, "", size=6)
    add_paragraph(tf,
        '"배열 주기가 줄어들수록 메타 원자 간 커플링에 의한 위상 왜곡이 어떻게 증가하는가?"',
        size=20, color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)

    # Sub-questions
    tb = add_textbox(slide, Inches(1.0), Inches(4.5), Inches(11.3), Inches(2.5))
    tf2 = tb.text_frame
    tf2.word_wrap = True
    set_text(tf2, "Sub-questions | 세부 질문", size=20, color=SUBTITLE_COLOR, bold=True)

    subs = [
        ("Isolated vs. array phase response — how different?",
         "고립 vs 배열 내 위상 응답의 차이는?"),
        ("Phase deviation trend as P varies from λ/2 to 2λ?",
         "P가 λ/2 → 2λ로 변할 때 위상 편차 추이는?"),
        ("Uniform vs. non-uniform array — different coupling?",
         "균일 배열 vs 비균일 배열에서 커플링 차이?"),
        ("Is there a threshold period where coupling becomes negligible?",
         "커플링 무시 가능한 주기 임계값 존재?"),
    ]
    for i, (en, kr) in enumerate(subs):
        add_paragraph(tf2, f"{i+1}. {en}", size=16, color=TEXT_COLOR)
        add_paragraph(tf2, f"    {kr}", size=14, color=MUTED_COLOR, space_before=Pt(0))


def slide_3_track(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 3, "Track A — Build a New Simulation", "시뮬레이션 직접 구축")

    shape = slide.shapes.add_shape(
        1, Inches(2.0), Inches(2.2), Inches(9.3), Inches(3.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = HIGHLIGHT_BG
    shape.line.color.rgb = TITLE_COLOR
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    set_text(tf, "Coupled Dipole Approximation (CDA)", size=28, color=ACCENT_COLOR, bold=True,
             alignment=PP_ALIGN.CENTER)

    add_paragraph(tf, "", size=10)
    add_bullet(tf, "Implement CDA simulation from scratch in Python",
               "CDA 시뮬레이션을 Python으로 처음부터 구현", size=20)
    add_paragraph(tf, "", size=6)
    add_bullet(tf, "Directly code Green's function + N×N linear system solver",
               "Green 함수 + N×N 선형 시스템 풀이를 직접 코딩", size=20)
    add_paragraph(tf, "", size=6)
    add_bullet(tf, "Python + NumPy / SciPy based",
               "Python + NumPy / SciPy 기반", size=20)


def slide_4_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 4, "Model", "물리 모델")

    y = Inches(1.5)
    left = Inches(0.8)

    # Dipole model
    tb1 = add_textbox(slide, left, y, Inches(5.5), Inches(0.4))
    set_text(tb1.text_frame, "Point dipole model | 점 쌍극자 모델:", size=18,
             color=SUBTITLE_COLOR, bold=True)
    y += Inches(0.45)
    add_equation_image(slide, r"p_i = \alpha_i \, E_{loc,i}", left + Inches(0.5), y)
    y += Inches(0.7)

    # Local field
    tb2 = add_textbox(slide, left, y, Inches(5.5), Inches(0.4))
    set_text(tb2.text_frame, "Local electric field | 국소 전기장:", size=18,
             color=SUBTITLE_COLOR, bold=True)
    y += Inches(0.45)
    add_equation_image(slide,
        r"E_{loc,i} = E_{inc,i} + \sum_{j \neq i} G(r_{ij}) \, p_j",
        left + Inches(0.5), y, fontsize=20)
    y += Inches(0.7)

    # Green's function
    tb3 = add_textbox(slide, left, y, Inches(5.5), Inches(0.4))
    set_text(tb3.text_frame, "Green's function (2D) | Green 함수:", size=18,
             color=SUBTITLE_COLOR, bold=True)
    y += Inches(0.45)
    add_equation_image(slide,
        r"G(r) = \frac{i}{4} H_0^{(1)}(k_0 \, r)",
        left + Inches(0.5), y)
    y += Inches(0.7)

    # Linear system
    tb4 = add_textbox(slide, left, y, Inches(5.5), Inches(0.4))
    set_text(tb4.text_frame, "Matrix form | 행렬 형태:", size=18,
             color=SUBTITLE_COLOR, bold=True)
    y += Inches(0.45)
    add_equation_image(slide, r"\mathbf{A}\,\mathbf{p} = \mathbf{E}_{inc}",
                       left + Inches(0.5), y, fontsize=24)

    # Right side: table-like info
    rt = Inches(7.0)
    tb5 = add_textbox(slide, rt, Inches(1.5), Inches(5.5), Inches(4.5))
    tf5 = tb5.text_frame
    tf5.word_wrap = True
    set_text(tf5, "Matrix element:", size=18, color=SUBTITLE_COLOR, bold=True)

    y_right = Inches(2.1)
    tb_eq = add_textbox(slide, rt + Inches(0.3), y_right, Inches(5.0), Inches(0.35))
    set_text(tb_eq.text_frame, "A_ii = 1/αᵢ  ,   A_ij = −G(rᵢⱼ)  (i≠j)",
             size=18, color=TEXT_COLOR)

    tb6 = add_textbox(slide, rt, Inches(3.3), Inches(5.5), Inches(3.5))
    tf6 = tb6.text_frame
    tf6.word_wrap = True
    set_text(tf6, "Key components | 주요 구성", size=18, color=SUBTITLE_COLOR, bold=True)
    add_paragraph(tf6, "", size=6)

    items = [
        ("State var | 상태 변수", "Dipole moment pᵢ (complex)\n쌍극자 모멘트 pᵢ (복소수)"),
        ("Parameters | 파라미터", "Period P, α (Lorentzian), λ\n주기 P, 분극률 α, 파장 λ"),
        ("Boundary | 경계 조건", "Plane wave, finite 1D array\n평면파, 유한 1D 배열"),
    ]
    for label, desc in items:
        add_paragraph(tf6, f"▸ {label}", size=16, color=ACCENT_COLOR, bold=True)
        add_paragraph(tf6, f"  {desc}", size=14, color=TEXT_COLOR, space_before=Pt(0))


def slide_5_io(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 5, "Input / Output", "입출력")

    # Input
    tb = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, "Input | 입력", size=22, color=ACCENT_COLOR, bold=True)

    inputs = [
        ("Wavelength λ", "파장 λ", "Design wavelength (normalized)"),
        ("Array period P", "배열 주기 P", "0.5λ ~ 2.0λ (sweep)"),
        ("Num. meta-atoms N", "메타 원자 수 N", "1D: 20 ~ 100"),
        ("Polarizability α", "분극률 α", "Lorentzian model"),
        ("Resonance freq. ω₀", "공진 주파수 ω₀", "May vary per atom"),
        ("Incident wave", "입사파", "Plane wave (normal incidence)"),
    ]
    for en, kr, val in inputs:
        add_paragraph(tf, f"▸ {en} | {kr}", size=16, color=SUBTITLE_COLOR, bold=True)
        add_paragraph(tf, f"   → {val}", size=14, color=TEXT_COLOR, space_before=Pt(0))

    # Output
    tb2 = add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(4.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    set_text(tf2, "Output | 출력", size=22, color=ACCENT_COLOR, bold=True)

    add_bullet(tf2, "Dipole moment pᵢ (complex)",
               "각 원자의 유도 쌍극자 모멘트 pᵢ (복소수)", size=16)
    add_paragraph(tf2, "", size=6)
    add_bullet(tf2, "Phase: arg(pᵢ) — actual phase in array",
               "pᵢ의 위상: 배열 내 실제 위상 응답", size=16)
    add_paragraph(tf2, "", size=6)
    add_bullet(tf2, "Isolated response: pᵢⁱˢᵒ = αᵢ · Eᵢₙ꜀",
               "고립 응답: pᵢⁱˢᵒ = αᵢ · Eᵢₙ꜀ (커플링 무시)", size=16)
    add_paragraph(tf2, "", size=6)

    tb3 = add_textbox(slide, Inches(7.0), Inches(4.5), Inches(5.5), Inches(1.5))
    set_text(tb3.text_frame, "Key output | 핵심 출력:", size=18, color=ACCENT_COLOR, bold=True)
    y_eq = Inches(5.0)
    add_equation_image(slide,
        r"\Delta\varphi_i = \arg(p_i^{coupled}) - \arg(p_i^{isolated})",
        Inches(7.3), y_eq, fontsize=20)


def slide_6_metric(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 6, "Quantitative Metrics", "측정 지표")

    left = Inches(0.8)
    y = Inches(1.6)

    metrics = [
        ("1. Mean Phase Deviation | 평균 위상 편차",
         r"\overline{\Delta\varphi} = \frac{1}{N}\sum_{i=1}^{N}|\Delta\varphi_i|"),
        ("2. Max Phase Deviation | 최대 위상 편차",
         r"\Delta\varphi_{\max} = \max_i |\Delta\varphi_i|"),
        ("3. Amplitude Distortion Ratio | 진폭 왜곡 비",
         r"ADR_i = \frac{|p_i^{coupled}|}{|p_i^{isolated}|}"),
        ("4. Optical Theorem Check | 에너지 보존 검증",
         r"\sigma_{ext} = \frac{4\pi k_0}{|E_0|^2} \sum_i Im(E_{inc,i}^* \cdot p_i)"),
    ]

    for title, eq in metrics:
        tb = add_textbox(slide, left, y, Inches(11.5), Inches(0.4))
        set_text(tb.text_frame, title, size=18, color=ACCENT_COLOR, bold=True)
        y += Inches(0.45)
        add_equation_image(slide, eq, left + Inches(0.5), y, fontsize=20)
        y += Inches(0.8)

    # note
    tb_note = add_textbox(slide, left, y + Inches(0.1), Inches(11.5), Inches(0.5))
    set_text(tb_note.text_frame,
             "Primary metric: plot  Δφ̄(P)  vs period P → quantify coupling effect",
             size=16, color=MUTED_COLOR)
    add_paragraph(tb_note.text_frame,
                  "주요 지표: Δφ̄(P) vs 주기 P 그래프 → 커플링 영향 정량화",
                  size=14, color=MUTED_COLOR, space_before=Pt(0))


def slide_7_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 7, "Controlled Comparison", "비교 실험 설계")

    # Main comparison table
    tb = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5))
    set_text(tb.text_frame, "Main: Period Sweep | 메인: 주기 스윕",
             size=20, color=ACCENT_COLOR, bold=True)

    rows = [
        ("Dense | 밀집", "0.5", "Strong | 강함", "Large distortion | 큰 왜곡"),
        ("Medium | 중간", "1.0", "Moderate | 중간", "Measurable | 측정 가능"),
        ("Sparse | 희박", "1.5", "Weak | 약함", "Small | 작음"),
        ("Very sparse | 매우 희박", "2.0", "Negligible | 무시 가능", "Baseline | 기준선"),
    ]

    y_start = Inches(2.2)
    headers = ["Period", "P/λ", "Coupling", "Expected Distortion"]
    x_positions = [Inches(0.8), Inches(3.5), Inches(5.5), Inches(8.0)]
    widths = [Inches(2.7), Inches(2.0), Inches(2.5), Inches(4.0)]

    for i, h in enumerate(headers):
        t = add_textbox(slide, x_positions[i], y_start, widths[i], Inches(0.35))
        set_text(t.text_frame, h, size=15, color=TITLE_COLOR, bold=True)

    for ri, row in enumerate(rows):
        y = y_start + Inches(0.4) * (ri + 1)
        for ci, cell in enumerate(row):
            t = add_textbox(slide, x_positions[ci], y, widths[ci], Inches(0.35))
            set_text(t.text_frame, cell, size=14, color=TEXT_COLOR)

    # Additional comparison
    tb2 = add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5))
    set_text(tb2.text_frame, "Additional: Uniform vs Non-uniform Array | 추가: 균일 vs 비균일 배열",
             size=20, color=ACCENT_COLOR, bold=True)

    tb3 = add_textbox(slide, Inches(0.8), Inches(5.2), Inches(5.5), Inches(1.8))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    set_text(tf3, "▸ Uniform array | 균일 배열", size=18, color=SUBTITLE_COLOR, bold=True)
    add_paragraph(tf3, "All atoms identical (same α)", size=15, color=TEXT_COLOR)
    add_paragraph(tf3, "모든 메타 원자 동일 → 순수 격자 커플링 분리", size=14, color=MUTED_COLOR,
                  space_before=Pt(0))

    tb4 = add_textbox(slide, Inches(6.8), Inches(5.2), Inches(5.5), Inches(1.8))
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    set_text(tf4, "▸ Non-uniform array | 비균일 배열", size=18, color=SUBTITLE_COLOR, bold=True)
    add_paragraph(tf4, "Each atom has different α (phase gradient)", size=15, color=TEXT_COLOR)
    add_paragraph(tf4, "각 원자 α 상이 → 실제 메타렌즈 상황과 유사", size=14, color=MUTED_COLOR,
                  space_before=Pt(0))


def slide_8_verification(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 8, "Verification Plan", "검증 계획")

    checks = [
        ("1. Analytical Comparison | 해석해 비교",
         "Sanity Check",
         ["Single atom: CDA result vs p = α · E_inc\n1개 원자: CDA 결과 vs 직접 계산",
          "Two atoms: compare with analytical 2-dipole solution\n2개 원자: 해석적 해와 비교"]),
        ("2. Limiting Case | 극한 검증",
         "Consistency Check",
         ["P → ∞ : coupled phase → isolated phase\nP → ∞ : 커플링 위상 → 고립 위상 수렴",
          "Δφ → 0  as  P → ∞"]),
        ("3. Optical Theorem | 에너지 보존",
         "Conservation Check",
         ["Verify extinction cross-section satisfies optical theorem\n소광 단면적이 optical theorem 만족하는지 검증"]),
        ("4. Convergence | 수렴성 검사",
         "Robustness Check",
         ["Central atom's Δφ converges as N increases\n배열 크기 N ↑ 시 중심 원자 Δφ 수렴 확인"]),
    ]

    left_col = Inches(0.6)
    right_col = Inches(6.8)
    y_positions = [Inches(1.6), Inches(3.6), Inches(1.6), Inches(3.6)]
    x_cols = [left_col, left_col, right_col, right_col]

    for i, (title, tag, items) in enumerate(checks):
        x = x_cols[i]
        y = y_positions[i]
        w = Inches(5.8)

        shape = slide.shapes.add_shape(1, x, y, w, Inches(1.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = HIGHLIGHT_BG
        shape.line.color.rgb = RGBColor(0x30, 0x40, 0x60)
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.word_wrap = True
        set_text(tf, title, size=16, color=ACCENT_COLOR, bold=True)
        add_paragraph(tf, f"[{tag}]", size=12, color=MUTED_COLOR, space_before=Pt(0))
        for item in items:
            add_paragraph(tf, f"  • {item}", size=13, color=TEXT_COLOR)


def slide_9_success(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_header(slide, 9, "Success Criteria", "성공 기준")

    # Minimum
    shape1 = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.5)
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = HIGHLIGHT_BG
    shape1.line.color.rgb = ACCENT_COLOR
    shape1.line.width = Pt(2)

    tf1 = shape1.text_frame
    tf1.word_wrap = True
    set_text(tf1, "Minimum Viable Project", size=22, color=ACCENT_COLOR, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_paragraph(tf1, "최소 완성 기준", size=16, color=MUTED_COLOR,
                  alignment=PP_ALIGN.CENTER, space_before=Pt(0))

    mins = [
        ("Build & solve CDA linear system",
         "CDA 선형 시스템 구축 및 풀이"),
        ("Extract phase deviation: isolated vs coupled",
         "고립 vs 커플링 위상 편차 추출"),
        ("Δφ̄(P) graph over period sweep",
         "주기 스윕에 따른 Δφ̄(P) 그래프"),
        ("Verify: single dipole matches analytical solution",
         "검증: 1개 쌍극자 해석해 일치 확인"),
    ]
    for i, (en, kr) in enumerate(mins):
        add_paragraph(tf1, f" {i+1}. {en}", size=16, color=TEXT_COLOR, bold=True)
        add_paragraph(tf1, f"     {kr}", size=14, color=MUTED_COLOR, space_before=Pt(0))

    # Full
    shape2 = slide.shapes.add_shape(
        1, Inches(7.0), Inches(1.6), Inches(5.5), Inches(4.5)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = HIGHLIGHT_BG
    shape2.line.color.rgb = TITLE_COLOR
    shape2.line.width = Pt(2)

    tf2 = shape2.text_frame
    tf2.word_wrap = True
    set_text(tf2, "Full Project (Target)", size=22, color=TITLE_COLOR, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_paragraph(tf2, "완전 프로젝트 (목표)", size=16, color=MUTED_COLOR,
                  alignment=PP_ALIGN.CENTER, space_before=Pt(0))

    fulls = [
        ("All minimum items +", "최소 항목 전부 +"),
        ("Uniform vs non-uniform comparison", "균일 vs 비균일 배열 비교"),
        ("Optical theorem verification", "Optical theorem 검증"),
        ("Array size convergence analysis", "배열 크기 수렴성 분석"),
        ("Quantify threshold period", "커플링 무시 가능 임계 주기 제시"),
        ("Physical interpretation & limitations", "물리적 해석 및 한계점 논의"),
    ]
    for en, kr in fulls:
        add_paragraph(tf2, f" ▸ {en}", size=16, color=TEXT_COLOR)
        add_paragraph(tf2, f"    {kr}", size=14, color=MUTED_COLOR, space_before=Pt(0))


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    make_title_slide(prs)      # Title
    slide_1_context(prs)       # 1. Research Context
    slide_2_question(prs)      # 2. Research Question
    slide_3_track(prs)         # 3. Track A
    slide_4_model(prs)         # 4. Model
    slide_5_io(prs)            # 5. Input / Output
    slide_6_metric(prs)        # 6. Metrics
    slide_7_comparison(prs)    # 7. Controlled Comparison
    slide_8_verification(prs)  # 8. Verification
    slide_9_success(prs)       # 9. Success Criteria

    prs.save(PPTX_PATH)
    print(f"PPT saved: {PPTX_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
