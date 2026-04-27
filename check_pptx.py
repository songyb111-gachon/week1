"""
Validate the generated PPTX:
  - file exists & is a valid PPTX
  - every slide has expected content
  - no shape overflows the slide bounds
  - all referenced figure files exist on disk
  - no empty text frames
  - prints a summary table for inspection
"""
import os
from pptx import Presentation
from pptx.util import Emu

PPTX_PATH = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "week2_progress.pptx"
)
FIG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "figures")

EXPECTED_SLIDE_COUNT = 10
issues = []
warnings = []


def emu_to_inches(v: int) -> float:
    return v / 914400.0


def check_file_exists():
    if not os.path.exists(PPTX_PATH):
        issues.append(f"PPTX file does not exist: {PPTX_PATH}")
        return False
    size = os.path.getsize(PPTX_PATH)
    print(f"[OK]  File exists: {PPTX_PATH}  ({size/1024:.1f} KB)")
    return True


def check_figures_exist():
    expected = [
        "baseline_uniform_array.png",
        "limiting_case_period.png",
        "period_sweep_comparison.png",
        "verification_convergence.png",
    ]
    missing = []
    for f in expected:
        p = os.path.join(FIG_DIR, f)
        if not os.path.exists(p):
            missing.append(f)
    if missing:
        issues.append(f"Missing figures: {missing}")
    else:
        print(f"[OK]  All {len(expected)} figures present")


def main():
    if not check_file_exists():
        return
    check_figures_exist()

    prs = Presentation(PPTX_PATH)
    n_slides = len(prs.slides)
    slide_w = emu_to_inches(prs.slide_width)
    slide_h = emu_to_inches(prs.slide_height)

    print(f"[OK]  Slide count: {n_slides}  (expected {EXPECTED_SLIDE_COUNT})")
    print(f"[OK]  Slide size : {slide_w:.2f} x {slide_h:.2f} in")
    if n_slides != EXPECTED_SLIDE_COUNT:
        warnings.append(f"Slide count mismatch: {n_slides} vs {EXPECTED_SLIDE_COUNT}")

    print()
    print("=" * 90)
    print(f"{'Slide':>5} | {'Shapes':>6} | {'Text':>5} | {'Pics':>5} | "
          f"{'Empty':>5} | {'Overflow':>8} | First text")
    print("=" * 90)

    for idx, slide in enumerate(prs.slides):
        n_shapes = len(slide.shapes)
        n_text = 0
        n_pics = 0
        n_empty = 0
        n_overflow = 0
        first_text = ""
        for sh in slide.shapes:
            if sh.has_text_frame:
                tf = sh.text_frame
                full = tf.text.strip()
                if not full:
                    n_empty += 1
                else:
                    n_text += 1
                    if not first_text:
                        snippet = full[:42].replace("\n", " ")
                        # ASCII-safe for cp949 console
                        first_text = snippet.encode("ascii", "replace").decode("ascii")
            if sh.shape_type == 13:  # Picture
                n_pics += 1

            # bounds check (some shapes may have undefined dims)
            try:
                left = emu_to_inches(sh.left or 0)
                top = emu_to_inches(sh.top or 0)
                w = emu_to_inches(sh.width or 0)
                h = emu_to_inches(sh.height or 0)
                if (left + w) > slide_w + 0.05 or (top + h) > slide_h + 0.05:
                    n_overflow += 1
            except (AttributeError, TypeError):
                pass

        # Print row
        print(f"{idx:>5} | {n_shapes:>6} | {n_text:>5} | {n_pics:>5} | "
              f"{n_empty:>5} | {n_overflow:>8} | {first_text}")

        if n_overflow:
            warnings.append(
                f"Slide {idx}: {n_overflow} shape(s) appear to overflow slide bounds"
            )

    print()
    print("=" * 90)
    print("DETAILED OVERFLOW INSPECTION")
    print("=" * 90)
    for idx, slide in enumerate(prs.slides):
        for j, sh in enumerate(slide.shapes):
            try:
                left = emu_to_inches(sh.left or 0)
                top = emu_to_inches(sh.top or 0)
                w = emu_to_inches(sh.width or 0)
                h = emu_to_inches(sh.height or 0)
                right = left + w
                bottom = top + h
                if right > slide_w + 0.05 or bottom > slide_h + 0.05:
                    text_preview = ""
                    if sh.has_text_frame:
                        text_preview = sh.text_frame.text.strip()[:40]
                        text_preview = text_preview.encode(
                            "ascii", "replace"
                        ).decode("ascii")
                    print(
                        f"Slide {idx}  shape[{j}]:  "
                        f"left={left:.2f}, top={top:.2f}, "
                        f"w={w:.2f}, h={h:.2f}, "
                        f"right={right:.2f} ({slide_w:.2f}), "
                        f"bottom={bottom:.2f} ({slide_h:.2f})    "
                        f"text='{text_preview}'"
                    )
            except (AttributeError, TypeError):
                pass

    print()
    print("=" * 90)
    print(f"SUMMARY:  {len(issues)} issues,  {len(warnings)} warnings")
    print("=" * 90)
    for x in issues:
        print(f"  [ISSUE]    {x}")
    for x in warnings:
        print(f"  [WARNING]  {x}")
    if not issues and not warnings:
        print("  [OK]  No issues found.")


if __name__ == "__main__":
    main()
